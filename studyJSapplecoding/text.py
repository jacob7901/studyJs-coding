from pptx import Presentation
from pptx.util import Inches

# 메뉴 트리 데이터
menu_tree = [
    ("공통(Common)", [
        ("로그인(Login)", [
            "로그인 가능(Support login Function)",
            "회원 가입 기능(Register)",
            "회원 정보 찾기 가능(Forgot password)",
        ]),
        ("화면(Screen)", [
            "반응형 웹 - PC(responsive)",
            "반응형 웹 - Tablet(responsive)",
            "반응형 웹 - Mobile (responsive)",
        ]),
        ("언어(Language)", [
            "Language - Korean",
            "Language - English",
            "Language - Vietnamese",
        ]),
    ]),
    ("매니저(Manager)", [
        "대시보드(Dashboard)",
        "수업관리",
        "강사 관리",
        "학생관리",
        "정산 관리",
        "메시지 관리",
    ]),
    ("강사(Tutor)", [
        "대시보드(Dashboard)",
        "프로필수정",
        "수업 관리",
    ]),
    ("학생(Student)", [
        "대시보드(Dashboard)",
        "프로필 수정",
        "수업 확인",
        "수업 취소",
    ]),
]

# PPT 프레젠테이션 생성
presentation = Presentation()

# 메뉴 트리 데이터를 PPT 슬라이드로 변환
for main_title, sub_menus in menu_tree:
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title = slide.shapes.title
    title.text = main_title

    if isinstance(sub_menus, list):
        for idx, sub_menu in enumerate(sub_menus):
            if isinstance(sub_menu, tuple):
                sub_title, sub_sub_menus = sub_menu
                left = Inches(1)
                top = Inches(1.5 + idx * 0.5)
                width = Inches(2)
                height = Inches(0.5)
                textbox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                text_frame.text = sub_title

                for sub_idx, sub_sub_menu in enumerate(sub_sub_menus):
                    left = Inches(3)
                    top = Inches(1.5 + idx * 0.5 + sub_idx * 0.5)
                    width = Inches(4)
                    height = Inches(0.5)
                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = textbox.text_frame
                    text_frame.text = sub_sub_menu
            else:
                left = Inches(1)
                top = Inches(1.5 + idx * 0.5)
                width = Inches(6)
                height = Inches(0.5)
                textbox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                text_frame.text = sub_menu

# PPT 파일 저장
presentation.save("menu_tree.pptx")
